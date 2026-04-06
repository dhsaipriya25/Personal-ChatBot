import os
import base64
import io
from flask import Flask, render_template, request, jsonify
import anthropic
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))


def extract_word_text(file_bytes, filename):
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
    return f"[Word Document: {filename}]\n\n{text}"


def extract_ppt_text(file_bytes, filename):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return f"[PowerPoint: {filename}]\n\n" + '\n'.join(lines)


def extract_excel_text(file_bytes, filename):
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    lines = [f"[Excel File: {filename}]"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"\n--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(cell.strip() for cell in row_data):
                lines.append(" | ".join(row_data))
    return '\n'.join(lines)


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/chat", methods=["POST"])
def chat():
    data = request.get_json()
    messages = data.get("messages", [])
    file_data = data.get("file")

    if not messages:
        return jsonify({"error": "No messages provided"}), 400

    if file_data:
        file_name = file_data.get("name", "file")
        file_type = file_data.get("type", "")
        file_b64 = file_data.get("data", "")
        file_bytes = base64.b64decode(file_b64)

        last_msg = messages[-1]
        user_text = last_msg.get("content", "") if isinstance(last_msg.get("content"), str) else ""

        if file_type.startswith("image/"):
            content = []
            if user_text:
                content.append({"type": "text", "text": user_text})
            content.append({
                "type": "image",
                "source": {"type": "base64", "media_type": file_type, "data": file_b64}
            })
            messages[-1] = {"role": "user", "content": content}

        elif file_type == "application/pdf":
            content = []
            if user_text:
                content.append({"type": "text", "text": user_text})
            content.append({
                "type": "document",
                "source": {"type": "base64", "media_type": "application/pdf", "data": file_b64}
            })
            messages[-1] = {"role": "user", "content": content}

        elif file_name.lower().endswith(".docx"):
            extracted = extract_word_text(file_bytes, file_name)
            combined = f"{extracted}\n\n{user_text}" if user_text else extracted
            messages[-1] = {"role": "user", "content": combined}

        elif file_name.lower().endswith(".pptx"):
            extracted = extract_ppt_text(file_bytes, file_name)
            combined = f"{extracted}\n\n{user_text}" if user_text else extracted
            messages[-1] = {"role": "user", "content": combined}

        elif file_name.lower().endswith((".xlsx", ".xls")):
            extracted = extract_excel_text(file_bytes, file_name)
            combined = f"{extracted}\n\n{user_text}" if user_text else extracted
            messages[-1] = {"role": "user", "content": combined}

    response = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=16000,
        system=(
            "You are a helpful assistant. When given files or documents, analyze them thoroughly "
            "and answer questions about their content. "
            "Always respond in plain text only. Use a dot (.) or dash (-) as bullet point indicators. "
            "Never use # for headings or * for bold/bullets. Write headings as plain text followed by a colon."
        ),
        messages=messages,
    )

    reply = next((b.text for b in response.content if b.type == "text"), "")
    return jsonify({"reply": reply})


if __name__ == "__main__":
    app.run(debug=True)
